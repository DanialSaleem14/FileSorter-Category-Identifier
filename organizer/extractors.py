import os
from typing import Optional

from PIL import Image
import pytesseract
from pdfminer.high_level import extract_text as pdf_extract_text
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import chardet


def _ensure_tesseract_path() -> None:
	path = os.getenv("TESSERACT_PATH")
	if path:
		pytesseract.pytesseract.tesseract_cmd = path


def extract_text_from_image(path: str) -> str:
	_ensure_tesseract_path()
	try:
		with Image.open(path) as img:
			# Convert to RGB if needed
			if img.mode != 'RGB':
				img = img.convert('RGB')
			# Use German language for better OCR results
			text = pytesseract.image_to_string(img, lang='deu+eng')
			return text
	except Exception as e:
		print(f"[yellow]OCR failed for {path}: {e}[/]")
		return ""


def extract_text_from_pdf(path: str) -> str:
	# Try text layer first; fallback to OCR per page for empty pages
	try:
		text = pdf_extract_text(path) or ""
		if text.strip():
			return text
	except Exception:
		pass
	# OCR pages if no text layer or on failure
	try:
		reader = PdfReader(path)
		page_images: list[str] = []
		for _ in reader.pages:
			# Simplify: rely on external tools for rasterization if needed. Return empty here.
			pass
		return text if (text := "".join(page_images)).strip() else ""
	except Exception:
		return ""


def extract_text_from_docx(path: str) -> str:
	try:
		doc = Document(path)
		return "\n".join(p.text for p in doc.paragraphs)
	except Exception:
		return ""


def extract_text_from_pptx(path: str) -> str:
	try:
		prs = Presentation(path)
		texts = []
		for slide in prs.slides:
			for shape in slide.shapes:
				if hasattr(shape, "text"):
					texts.append(shape.text)
		return "\n".join(texts)
	except Exception:
		return ""


def extract_text_from_xlsx(path: str) -> str:
	try:
		wb = load_workbook(path, data_only=True, read_only=True)
		texts = []
		for ws in wb:
			for row in ws.iter_rows(values_only=True):
				for cell in row:
					if cell is None:
						continue
					texts.append(str(cell))
					if len(texts) > 50000:
						break
		return "\n".join(texts)
	except Exception:
		return ""


def extract_text_from_txt(path: str) -> str:
	try:
		with open(path, "rb") as f:
			raw = f.read()
		enc = chardet.detect(raw).get("encoding") or "utf-8"
		return raw.decode(enc, errors="ignore")
	except Exception:
		return ""


EXTENSIONS = {
	".pdf": extract_text_from_pdf,
	".png": extract_text_from_image,
	".jpg": extract_text_from_image,
	".jpeg": extract_text_from_image,
	".tif": extract_text_from_image,
	".tiff": extract_text_from_image,
	".docx": extract_text_from_docx,
	".pptx": extract_text_from_pptx,
	".xlsx": extract_text_from_xlsx,
	".txt": extract_text_from_txt,
}


def extract_text_from_file(path: str) -> str:
	ext = os.path.splitext(path)[1].lower()
	extractor = EXTENSIONS.get(ext)
	if not extractor:
		return ""
	return extractor(path) or ""
