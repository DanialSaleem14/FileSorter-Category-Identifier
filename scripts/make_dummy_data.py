import os
from docx import Document
from pptx import Presentation
from openpyxl import Workbook

BASE = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
UNSORTED = os.path.join(BASE, "UnsortedDemo")

CATEGORIES_AND_TEXTS = [
	("Rechnungen", "Rechnung Nr. 1024 Betrag 199,00 EUR zzgl. MwSt."),
	("Mahnungen", "Zahlungserinnerung Letzte Mahnung fällig"),
	("Quittungen", "Quittung über den Betrag von 25,00 EUR"),
	("Angebote", "Preisangebot für IT-Dienstleistungen"),
	("Bestellungen", "Bestellung Auftragsbestätigung Nr. 77"),
	("Arbeitsvertrag", "Arbeitsvertrag zwischen Arbeitgeber und Arbeitnehmer"),
	("Mietvertrag", "Mietvertrag Kaution Mietobjekt"),
	("Kaufvertrag", "Kaufvertrag Käufer Verkäufer Kaufpreis"),
	("Kontoauszuege_Bank", "Kontoauszug IBAN DE12 3456 7890 1234 5678 90"),
	("Praesentationen", "Präsentation Agenda Projektziele"),
	("Zertifikate", "Zertifikat Projektschulung erfolgreich abgeschlossen"),
	("Rezepte", "Rezept Ibuprofen 400mg"),
]


def make_dirs() -> None:
	os.makedirs(UNSORTED, exist_ok=True)


def make_docx(name: str, text: str) -> None:
	d = Document()
	d.add_heading(name, 0)
	d.add_paragraph(text)
	d.save(os.path.join(UNSORTED, f"{name}.docx"))


def make_pptx(name: str, text: str) -> None:
	prs = Presentation()
	slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(slide_layout)
	slide.shapes.title.text = name
	slide.placeholders[1].text = text
	prs.save(os.path.join(UNSORTED, f"{name}.pptx"))


def make_xlsx(name: str, text: str) -> None:
	wb = Workbook()
	ws = wb.active
	ws["A1"] = name
	ws["A2"] = text
	wb.save(os.path.join(UNSORTED, f"{name}.xlsx"))


def main() -> None:
	make_dirs()
	for cat, text in CATEGORIES_AND_TEXTS:
		make_docx(cat, text)
		make_pptx(cat, text)
		make_xlsx(cat, text)
	with open(os.path.join(UNSORTED, "unbekannt.txt"), "w", encoding="utf-8") as f:
		f.write("Dies ist eine Notiz ohne klare Kategorie.")
	
	# Create test images
	import subprocess
	try:
		subprocess.run(["python", "scripts/create_test_images.py"], check=True)
	except:
		print("Could not create test images (PIL might not be available)")
	
	print(f"Dummy files created in: {UNSORTED}")


if __name__ == "__main__":
	main()
